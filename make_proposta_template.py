"""
make_proposta_template.py
Cria templates/proposta.pptx a partir do arquivo original.
Substitui os dados de exemplo por marcadores {{PLACEHOLDER}}.
Execute uma vez sempre que o design do original for atualizado.
"""
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC  = Path(r"C:\Users\prpau\OneDrive\Desktop\Propostas de Honorarios\proposta_A4.pptx")
DEST = Path(__file__).parent / "templates" / "proposta.pptx"

# Slide 1 — capa: info do solicitante principal
SLIDE1_REPL = [
    ("Anderson Hirata de Moura",         "{{SOLICITANTE_NOME}}"),
    ("Hirata Sistemas Rurais e TI Ltda", "{{SOLICITANTE_EMPRESA}}"),
    ("Execução de Título Extrajudicial", "{{ASSUNTO}}"),
    ("Abril de 2026",                    "{{DATA_PROPOSTA}}"),
]

# Slide 3 — partes envolvidas: lista dinâmica (pode ter N partes)
# {{PARTES_LISTADAS}} recebe todas as partes separadas por \n
# {{PARTES_S3_CLEAR}} é esvaziado (era o segundo parágrafo com o nome)
SLIDE3_REPL = [
    ("Hirata Sistemas Rurais e TI Ltda", "{{PARTES_LISTADAS}}"),
    ("Anderson Hirata de Moura",         "{{PARTES_S3_CLEAR}}"),
    ("Execução de Título Extrajudicial", "{{ASSUNTO}}"),
]

# Slides 5, 6, 7 — valores financeiros (ordem importa: mais longos primeiro)
FINANCIAL_REPL = [
    ("Economia de R$ 500,00 em relação ao parcelado por boleto. Aprovação rápida.", "{{OP2_ECONOMIA_LONGA}}"),
    ("Até 6 parcelas sem juros adicionais", "{{OP2_PARCELAS_SEM_JUROS}}"),
    ("Economia de R$ 500,00",              "{{OP2_ECONOMIA_CURTA}}"),
    ("Total: R$ 2.502,00",                 "{{OP2_TOTAL}}"),
    ("em 6 parcelas mensais",              "{{OP2_PARCELAS_MENSAIS}}"),
    ("5x R$ 500,00",                       "{{OP1_PARCELAS}}"),
    ("Até 6x de",                          "{{OP2_ATE_PARCELAS}}"),
    ("R$ 3.000,00",                        "{{OP1_VALOR_TOTAL}}"),
    ("R$ 2.500,00",                        "{{OP2_AVISTA}}"),
    ("R$ 500,00",                          "{{OP1_ENTRADA}}"),
    ("R$ 417,00",                          "{{OP2_VALOR_PARCELA}}"),
]

SLIDE_RULES = {
    1: SLIDE1_REPL,
    3: SLIDE3_REPL,
    5: FINANCIAL_REPL,
    6: FINANCIAL_REPL,
    7: FINANCIAL_REPL,
}


def _replace_para(para, old: str, new: str) -> bool:
    full = "".join(r.text or "" for r in para.runs)
    if not full or old not in full:
        return False
    new_full = full.replace(old, new)
    if para.runs:
        para.runs[0].text = new_full
        for r in para.runs[1:]:
            r.text = ""
    return True


def _process_shapes(shapes, repl_list: list, replaced: dict) -> None:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _process_shapes(shape.shapes, repl_list, replaced)
        elif shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for old, new in repl_list:
                    if _replace_para(para, old, new):
                        replaced[old] = replaced.get(old, 0) + 1


def main() -> None:
    if not SRC.exists():
        print(f"Arquivo original não encontrado: {SRC}")
        return

    prs = Presentation(str(SRC))
    replaced: dict[str, int] = {}

    for idx, slide in enumerate(prs.slides, 1):
        repl = SLIDE_RULES.get(idx, [])
        if repl:
            _process_shapes(slide.shapes, repl, replaced)

    DEST.parent.mkdir(exist_ok=True)
    prs.save(str(DEST))

    print(f"Template salvo em: {DEST}\n")
    print("Substituições realizadas:")
    for old, count in replaced.items():
        print(f"  [{count}x] {old!r}")

    all_patterns = set(old for rules in SLIDE_RULES.values() for old, _ in rules)
    missing = [p for p in all_patterns if p not in replaced]
    if missing:
        print("\nAVISO — padrões não encontrados:")
        for m in missing:
            print(f"  - {m!r}")


if __name__ == "__main__":
    main()
