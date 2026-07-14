"""
document_generator.py
Toda a lógica de geração de contratos, procuração e hipossuficiência.
Sem dependência de internet, Gemini ou Bitrix.
"""
from __future__ import annotations

import calendar
import re
import shutil
from datetime import date, datetime
from pathlib import Path
from typing import Any

from docxtpl import DocxTemplate, RichText

# ── RichText helpers ─────────────────────────────────────────────────────────

_BOLD_6   = {'CONTRATANTE', 'CONTRATADA'}
_RT_FONT  = 'Arial'
_RT_SIZE  = 20  # meios-pontos → 10pt

def _rt_bold(text: str, bold_words: set) -> RichText:
    """Cria RichText com as palavras indicadas em negrito, fonte e tamanho fixos."""
    rt = RichText()
    pattern = '(' + '|'.join(re.escape(w) for w in sorted(bold_words, key=len, reverse=True)) + ')'
    for i, line in enumerate(text.split('\n')):
        if i > 0:
            rt.add('\n')
        for part in re.split(pattern, line):
            if part:
                rt.add(part, bold=(part in bold_words), font=_RT_FONT, size=_RT_SIZE)
    return rt

# ── Meses em português ────────────────────────────────────────────────────────
_MESES = [
    "", "janeiro", "fevereiro", "março", "abril", "maio", "junho",
    "julho", "agosto", "setembro", "outubro", "novembro", "dezembro",
]

# ── Números por extenso ───────────────────────────────────────────────────────
_ONES = [
    "", "um", "dois", "três", "quatro", "cinco", "seis", "sete", "oito", "nove",
    "dez", "onze", "doze", "treze", "quatorze", "quinze", "dezesseis",
    "dezessete", "dezoito", "dezenove",
]
_TENS = ["", "", "vinte", "trinta", "quarenta", "cinquenta",
         "sessenta", "setenta", "oitenta", "noventa"]
_HUNDREDS = [
    "", "cento", "duzentos", "trezentos", "quatrocentos", "quinhentos",
    "seiscentos", "setecentos", "oitocentos", "novecentos",
]


def _int_extenso(n: int) -> str:
    if n == 0:   return "zero"
    if n == 100: return "cem"
    if n < 20:   return _ONES[n]
    if n < 100:
        d, u = divmod(n, 10)
        return _TENS[d] + (" e " + _ONES[u] if u else "")
    c, r = divmod(n, 100)
    return _HUNDREDS[c] + (" e " + _int_extenso(r) if r else "")


def valor_extenso(value_str: str) -> str:
    """Converte '6.000,00' → 'seis mil reais'."""
    try:
        clean = value_str.replace("R$", "").replace(".", "").replace(",", ".").strip()
        value = float(clean)
    except Exception:
        return value_str

    cents = round(value * 100)
    reais, centavos = divmod(cents, 100)

    parts: list[str] = []
    tmp = reais
    if tmp >= 1_000_000:
        m, tmp = divmod(tmp, 1_000_000)
        parts.append(_int_extenso(m) + (" milhão" if m == 1 else " milhões"))
    if tmp >= 1_000:
        t, tmp = divmod(tmp, 1_000)
        parts.append(_int_extenso(t) + " mil")
    if tmp:
        parts.append(_int_extenso(tmp))

    if not parts:
        text = "zero reais"
    else:
        noun = "real" if reais == 1 else "reais"
        joined = " e ".join(parts) if len(parts) <= 2 else ", ".join(parts[:-1]) + " e " + parts[-1]
        text = joined + " " + noun

    if centavos:
        text += " e " + _int_extenso(centavos) + " centavos"
    return text


def _fmt_brl(value_str: str) -> str:
    """Garante formatação 'R$ 1.000,00'."""
    s = value_str.strip().lstrip("R$").strip()
    if not s:
        return ""
    return f"R$ {s}"


def _add_months(d: date, n: int) -> date:
    month = d.month - 1 + n
    year  = d.year + month // 12
    month = month % 12 + 1
    day   = min(d.day, calendar.monthrange(year, month)[1])
    return date(year, month, day)


def _parse_date(s: str) -> date:
    """Tenta DD/MM/AAAA."""
    try:
        d, m, y = map(int, s.strip().split("/"))
        return date(y, m, d)
    except Exception:
        return date.today()


# ── Geração da Cláusula 6 ─────────────────────────────────────────────────────

def build_clausula_6(payment: dict) -> tuple[str, str]:
    """
    Retorna (texto_6_1, texto_6_2).
    texto_6_2 é string vazia quando não deve existir.
    """
    tipo          = payment.get("tipo", "boleto")
    valor_total   = payment.get("valor_total", "")
    total_brl     = _fmt_brl(valor_total)
    total_ext     = valor_extenso(valor_total)

    _bold = _BOLD_6 | {'6.1.', '6.1'}

    if tipo == "avista":
        data_pgto = payment.get("data", "")
        txt = (
            f"6.1. Pela prestação dos serviços descritos neste contrato, "
            f"O CONTRATANTE pagará à CONTRATADA o valor de {total_brl} "
            f"({total_ext}), de forma à vista na data de {data_pgto}, "
            f"conforme combinado entre as partes."
        )
        return _rt_bold(txt, _bold), ""

    if tipo == "cartao":
        n_parc    = payment.get("n_parcelas", "")
        val_parc  = _fmt_brl(payment.get("valor_parcela", ""))
        parc_ext  = valor_extenso(payment.get("valor_parcela", ""))
        n_ext     = _int_extenso(int(n_parc)) if n_parc.isdigit() else n_parc
        data_pgto = payment.get("data", "")
        txt = (
            f"6.1. Pela prestação dos serviços descritos neste contrato, "
            f"O CONTRATANTE pagará à CONTRATADA o valor de {total_brl} "
            f"({total_ext}), podendo o pagamento ser realizado à vista ou "
            f"parcelado no cartão de crédito em até {n_parc} ({n_ext}) "
            f"parcelas de {val_parc} ({parc_ext}), na data de {data_pgto}, "
            f"conforme condições acordadas entre as partes."
        )
        return _rt_bold(txt, _bold), ""

    # boleto / PIX
    valor_entrada = payment.get("valor_entrada", "")
    data_entrada  = payment.get("data_entrada", "")
    n_parc        = payment.get("n_parcelas", "")
    val_parc      = _fmt_brl(payment.get("valor_parcela", ""))
    entrada_brl   = _fmt_brl(valor_entrada)
    entrada_ext   = valor_extenso(valor_entrada)
    parc_ext      = valor_extenso(payment.get("valor_parcela", ""))
    n_int         = int(n_parc) if str(n_parc).isdigit() else 0
    n_ext         = _int_extenso(n_int)

    txt = (
        f"6.1 – Pela prestação dos serviços descritos neste contrato, "
        f"o CONTRATANTE pagará à CONTRATADA o valor de {total_brl} "
        f"({total_ext}), com entrada de {entrada_brl} ({entrada_ext}) "
        f"na data de {data_entrada}, e o saldo restante em {n_ext} "
        f"parcelas de {val_parc} ({parc_ext}) cada, a serem pagas nas "
        f"datas subsequentes, conforme combinado entre as partes:"
    )

    base = _parse_date(data_entrada)
    linhas: list[str] = []
    for i in range(1, n_int + 1):
        parc_date = _add_months(base, i)
        linhas.append(f"• {i + 1}ª parcela: {parc_date.strftime('%d/%m/%Y')} – no valor de {val_parc} ({parc_ext})")
    if linhas:
        txt += "\n" + "\n".join(linhas)

    c62 = (
        "6.2. O pagamento será realizado por PIX IDENTIFICADO NA CONTA "
        "BANCÁRIA DA CONTRATADA, com vencimento de cada mês."
    )
    return _rt_bold(txt, _bold), c62


# ── Quadro Resumo ─────────────────────────────────────────────────────────────

def build_quadro_custos(payment: dict) -> dict:
    """Retorna dict com campos para seção 3 do Quadro Resumo."""
    tipo = payment.get("tipo", "boleto")

    if tipo == "avista":
        return {
            "custos_desc":   f"R$ {payment.get('valor_total', '')} à vista.",
            "parcelamento":  "",
            "n_parcelas_qr": "",
        }
    if tipo == "cartao":
        vt   = payment.get("valor_total", "")
        np   = payment.get("n_parcelas", "")
        vp   = payment.get("valor_parcela", "")
        return {
            "custos_desc":   f"R$ {vt} em {np}x de R${vp}",
            "parcelamento":  "Crédito",
            "n_parcelas_qr": np,
        }
    # boleto
    ve = payment.get("valor_entrada", "")
    np = payment.get("n_parcelas", "")
    vp = payment.get("valor_parcela", "")
    n_total = int(np) + 1 if str(np).isdigit() else 0
    return {
        "custos_desc":   f"ENTRADA DE R${ve} + {np}X DE R${vp}",
        "parcelamento":  "Boleto",
        "n_parcelas_qr": str(n_total),
    }


# ── Identificação do cliente (parágrafo padrão) ───────────────────────────────

# ── Estado civil com concordância de gênero ───────────────────────────────────

_EC_MASC: dict[str, str] = {
    "solteiro":      "solteiro",
    "casado":        "casado",
    "divorciado":    "divorciado",
    "viuvo":         "viúvo",
    "separado":      "separado de fato",
    "uniao_estavel": "união estável",
}
_EC_FEM: dict[str, str] = {
    "solteiro":      "solteira",
    "casado":        "casada",
    "divorciado":    "divorciada",
    "viuvo":         "viúva",
    "separado":      "separada de fato",
    "uniao_estavel": "união estável",
}

def _ec_gendered(ec: str, fem: bool) -> str:
    d = _EC_FEM if fem else _EC_MASC
    return d.get(ec.lower().strip(), ec)


def build_bloco_hipo(client: dict) -> RichText:
    """Bloco do cliente para a hipossuficiência — Tahoma 12pt, sem ponto final."""
    fem       = client.get("nacionalidade", "").lower().endswith("a")
    inscrito  = "inscrita"                if fem else "inscrito"
    residente = "residente e domiciliada" if fem else "residente e domiciliado"
    ec        = _ec_gendered(client.get("estado_civil", ""), fem)

    rua    = client.get("rua", "")
    bairro = client.get("bairro", "")
    comp   = client.get("complemento", "")
    cidade = client.get("cidade", "")
    uf     = client.get("uf", "")
    cep    = client.get("cep", "")

    end_parts = [p for p in [rua, bairro, comp] if p]
    endereco  = ", ".join(end_parts)

    kw = dict(font="Tahoma", size=24)  # 12pt em meios-pontos
    rt = RichText()
    rt.add(client.get("nome", ""), bold=True, **kw)
    rt.add(
        f", {client.get('nacionalidade', '')}, {ec}, "
        f"{client.get('profissao', '')}, {inscrito} no CPF/MF sob o nº {client.get('cpf', '')}",
        **kw
    )
    rg = client.get("rg", "")
    if rg:
        rt.add(f", portador(a) da carteira de identidade nº {rg}", **kw)
    email = client.get("email", "")
    if email:
        rt.add(f", Email: {email}", **kw)
    rt.add(f", {residente} na {endereco}", **kw)
    if cidade and uf:
        rt.add(f", {cidade}/{uf}", **kw)
    if cep:
        rt.add(f", CEP: {cep}", **kw)
    return rt


def build_bloco_procuracao(client: dict) -> RichText:
    """Bloco do OUTORGANTE para a procuração — Tahoma 11pt."""
    fem       = client.get("nacionalidade", "").lower().endswith("a")
    inscrito  = "inscrita"                if fem else "inscrito"
    residente = "residente e domiciliada" if fem else "residente e domiciliado"
    ec        = _ec_gendered(client.get("estado_civil", ""), fem)

    rua    = client.get("rua", "")
    bairro = client.get("bairro", "")
    comp   = client.get("complemento", "")
    cidade = client.get("cidade", "")
    uf     = client.get("uf", "")
    cep    = client.get("cep", "")

    end_parts = [p for p in [rua, bairro, comp] if p]
    endereco  = ", ".join(end_parts)

    kw = dict(font="Tahoma", size=22)  # 11pt em meios-pontos
    rt = RichText()
    rt.add(client.get("nome", ""), bold=True, **kw)
    rt.add(
        f", {client.get('nacionalidade', '')}, {ec}, "
        f"{client.get('profissao', '')}, {inscrito} no CPF/MF sob o nº {client.get('cpf', '')}",
        **kw
    )
    rg = client.get("rg", "")
    if rg:
        rt.add(f", portador(a) da carteira de identidade nº {rg}", **kw)
    email = client.get("email", "")
    if email:
        rt.add(f", Email: {email}", **kw)
    rt.add(f", {residente} na {endereco}", **kw)
    if cidade and uf:
        rt.add(f", {cidade}/{uf}", **kw)
    if cep:
        rt.add(f", CEP: {cep}", **kw)
    rt.add(".", **kw)
    return rt


def build_bloco_procuracao_pj(client: dict, empresa: dict) -> RichText:
    """Bloco do OUTORGANTE PJ: empresa + representante (sócio)."""
    fem       = client.get("nacionalidade", "").lower().endswith("a")
    inscrito  = "inscrita" if fem else "inscrito"
    residente = "residente e domiciliada" if fem else "residente e domiciliado"
    ec        = _ec_gendered(client.get("estado_civil", ""), fem)

    # Endereço da empresa
    e_rua   = empresa.get("rua", "")
    e_bairro = empresa.get("bairro", "")
    e_comp  = empresa.get("complemento", "")
    e_cidade = empresa.get("cidade", "")
    e_uf    = empresa.get("uf", "")
    e_cep   = empresa.get("cep", "")
    e_end   = ", ".join(p for p in [e_rua, e_bairro, e_comp] if p)

    # Endereço do sócio
    rua     = client.get("rua", "")
    bairro  = client.get("bairro", "")
    comp    = client.get("complemento", "")
    cidade  = client.get("cidade", "")
    uf      = client.get("uf", "")
    cep     = client.get("cep", "")
    endereco = ", ".join(p for p in [rua, bairro, comp] if p)

    kw = dict(font="Tahoma", size=22)  # 11pt em meios-pontos
    rt = RichText()
    rt.add(empresa.get("nome", ""), bold=True, **kw)
    rt.add(
        f", pessoa jurídica de direito privado, inscrita no CNPJ sob o nº "
        f"{empresa.get('cnpj', '')}, com sede à {e_end}",
        **kw
    )
    if e_cidade and e_uf:
        rt.add(f", da cidade de {e_cidade}/{e_uf}", **kw)
    if e_cep:
        rt.add(f", CEP {e_cep}", **kw)
    rt.add(", neste ato representada por seu sócio ", **kw)
    rt.add(client.get("nome", ""), bold=True, **kw)
    rt.add(
        f", {client.get('nacionalidade', '')}, {ec}, "
        f"{client.get('profissao', '')}, {inscrito} no CPF/MF sob o nº {client.get('cpf', '')}",
        **kw
    )
    rg = client.get("rg", "")
    if rg:
        rt.add(f", portador(a) da carteira de identidade nº {rg}", **kw)
    email = client.get("email", "")
    if email:
        rt.add(f", Email: {email}", **kw)
    rt.add(f", {residente} na {endereco}", **kw)
    if cidade and uf:
        rt.add(f", {cidade}/{uf}", **kw)
    if cep:
        rt.add(f", CEP: {cep}", **kw)
    rt.add(".", **kw)
    return rt


def build_bloco_cliente(client: dict) -> RichText:
    fem       = client.get("nacionalidade", "").lower().endswith("a")
    inscrito  = "inscrita"                if fem else "inscrito"
    residente = "residente e domiciliada" if fem else "residente e domiciliado"
    ec        = _ec_gendered(client.get("estado_civil", ""), fem)

    rua    = client.get("rua", "")
    bairro = client.get("bairro", "")
    comp   = client.get("complemento", "")
    cidade = client.get("cidade", "")
    uf     = client.get("uf", "")
    cep    = client.get("cep", "")

    end_parts = [p for p in [rua, bairro, comp] if p]
    endereco  = ", ".join(end_parts)

    kw = dict(font=_RT_FONT, size=_RT_SIZE)
    rt = RichText()
    rt.add(client.get("nome", ""), bold=True, **kw)
    rt.add(
        f", {client.get('nacionalidade', '')}, {ec}, "
        f"{client.get('profissao', '')}, {inscrito} no CPF/MF sob o nº {client.get('cpf', '')}",
        **kw
    )
    rg = client.get("rg", "")
    if rg:
        rt.add(f", portador(a) da carteira de identidade nº {rg}", **kw)
    email = client.get("email", "")
    if email:
        rt.add(f", Email: {email}", **kw)
    telefone = client.get("telefone", "")
    if telefone:
        rt.add(f", tel: {telefone}", **kw)
    rt.add(f", {residente} na {endereco}", **kw)
    if cidade and uf:
        rt.add(f", {cidade}/{uf}", **kw)
    if cep:
        rt.add(f", CEP: {cep}", **kw)
    rt.add(', doravante denominado ', **kw)
    rt.add('"Contratante"', bold=True, **kw)
    rt.add(';', **kw)
    return rt


def build_bloco_cliente_pj(client: dict, empresa: dict) -> RichText:
    """Bloco do CONTRATANTE PJ para o contrato — Arial 10pt."""
    fem       = client.get("nacionalidade", "").lower().endswith("a")
    inscrito  = "inscrita" if fem else "inscrito"
    residente = "residente e domiciliada" if fem else "residente e domiciliado"
    ec        = _ec_gendered(client.get("estado_civil", ""), fem)

    e_rua    = empresa.get("rua", "")
    e_bairro = empresa.get("bairro", "")
    e_comp   = empresa.get("complemento", "")
    e_cidade = empresa.get("cidade", "")
    e_uf     = empresa.get("uf", "")
    e_cep    = empresa.get("cep", "")
    e_end    = ", ".join(p for p in [e_rua, e_bairro, e_comp] if p)

    rua      = client.get("rua", "")
    bairro   = client.get("bairro", "")
    comp     = client.get("complemento", "")
    cidade   = client.get("cidade", "")
    uf       = client.get("uf", "")
    cep      = client.get("cep", "")
    endereco = ", ".join(p for p in [rua, bairro, comp] if p)

    kw = dict(font=_RT_FONT, size=_RT_SIZE)
    rt = RichText()
    rt.add(empresa.get("nome", ""), bold=True, **kw)
    rt.add(
        f", pessoa jurídica de direito privado, inscrita no CNPJ sob o nº "
        f"{empresa.get('cnpj', '')}, com sede à {e_end}",
        **kw
    )
    if e_cidade and e_uf:
        rt.add(f", da cidade de {e_cidade}/{e_uf}", **kw)
    if e_cep:
        rt.add(f", CEP {e_cep}", **kw)
    rt.add(", neste ato representada por seu sócio ", **kw)
    rt.add(client.get("nome", ""), bold=True, **kw)
    rt.add(
        f", {client.get('nacionalidade', '')}, {ec}, "
        f"{client.get('profissao', '')}, {inscrito} no CPF/MF sob o nº {client.get('cpf', '')}",
        **kw
    )
    rg = client.get("rg", "")
    if rg:
        rt.add(f", portador(a) da carteira de identidade nº {rg}", **kw)
    email = client.get("email", "")
    if email:
        rt.add(f", Email: {email}", **kw)
    telefone = client.get("telefone", "")
    if telefone:
        rt.add(f", tel: {telefone}", **kw)
    rt.add(f", {residente} na {endereco}", **kw)
    if cidade and uf:
        rt.add(f", {cidade}/{uf}", **kw)
    if cep:
        rt.add(f", CEP: {cep}", **kw)
    rt.add(', doravante denominado ', **kw)
    rt.add('"Contratante"', bold=True, **kw)
    rt.add(';', **kw)
    return rt


# ── Data por extenso ──────────────────────────────────────────────────────────

def build_local_data(cidade: str, uf: str = "") -> str:
    today = date.today()
    loc   = cidade.title() if cidade else (uf or "")
    return f"{loc}, {today.day:02d} de {_MESES[today.month]} de {today.year}."


_OBS_FORA_ENDERECO = "O VEÍCULO NÃO ESTÁ LOCALIZADO NO MESMO ENDEREÇO DO CARNÊ DE FINANCIAMENTO"

def _build_veiculo_obs(vehicle: dict) -> str:
    parts = []
    if vehicle.get("fora_endereco", False):
        parts.append(_OBS_FORA_ENDERECO)
    extra = vehicle.get("observacao", "").strip()
    if extra:
        parts.append(extra)
    return " - ".join(parts)


# ── Contexto completo do template ─────────────────────────────────────────────

def build_context(payload: dict) -> dict[str, Any]:
    client       = payload.get("client", {})
    contract_type = payload.get("contractType", "emprestimo")
    processes    = payload.get("processes", [{}])
    vehicle      = payload.get("vehicle") or {}
    payment      = payload.get("payment", {})
    empresa      = payload.get("empresa") or {}

    c61, c62 = build_clausula_6(payment)
    qr_custos = build_quadro_custos(payment)

    # Taxa de rescisão = valor da entrada (boleto) ou valor total (cartão/à vista)
    tipo_pgto = payment.get("tipo", "boleto")
    if tipo_pgto == "boleto":
        _taxa_str = payment.get("valor_entrada", "")
    else:
        _taxa_str = payment.get("valor_total", "")
    taxa_brl    = _fmt_brl(_taxa_str)
    taxa_extenso = valor_extenso(_taxa_str)

    # Processo(s) — pega o primeiro para campos simples, lista completa para loop
    proc0 = processes[0] if processes else {}

    def proc_existe(p: dict) -> str:
        ep  = p.get("existe_processo", "nao")
        num = p.get("numero_processo", "")
        if ep == "nao":
            return "Não"
        return f"Sim, {num}" if num else "Sim"

    processos_ctx = [
        {
            "banco":             p.get("banco", ""),
            "divida":            _fmt_brl(p.get("divida", "")),
            "valor_parcela":     _fmt_brl(p.get("valor_parcela", "")),
            "parcelas_totais":   p.get("parcelas_totais", ""),
            "parcelas_pagas":    p.get("parcelas_pagas", ""),
            "parcelas_abertas":  p.get("parcelas_abertas", ""),
            "parcelas_vencidas": p.get("parcelas_vencidas", ""),
            "existe_processo":   proc_existe(p),
        }
        for p in processes
    ]

    ctx: dict[str, Any] = {
        # Cliente
        "nome":           client.get("nome", ""),
        "nacionalidade":  client.get("nacionalidade", ""),
        "estado_civil":   client.get("estado_civil", ""),
        "profissao":      client.get("profissao", ""),
        "cpf":            client.get("cpf", ""),
        "rg":             client.get("rg", ""),
        "email":          client.get("email", ""),
        "telefone":       client.get("telefone", ""),
        "rua":            client.get("rua", ""),
        "bairro":         client.get("bairro", ""),
        "complemento":    client.get("complemento", ""),
        "cidade":         client.get("cidade", ""),
        "uf":             client.get("uf", ""),
        "cep":            client.get("cep", ""),
        "bloco_cliente":     build_bloco_cliente(client),
        "bloco_procuracao":     build_bloco_procuracao(client),
        "bloco_procuracao_pj":  build_bloco_procuracao_pj(client, empresa),
        "bloco_hipo":        build_bloco_hipo(client),
        "local_data":        build_local_data(client.get("cidade", ""), client.get("uf", "")),

        # Processo 1 (atalho)
        "banco":             proc0.get("banco", ""),
        "divida":            _fmt_brl(proc0.get("divida", "")),
        "valor_parcela_banco": _fmt_brl(proc0.get("valor_parcela", "")),
        "parcelas_totais":   proc0.get("parcelas_totais", ""),
        "parcelas_pagas":    proc0.get("parcelas_pagas", ""),
        "parcelas_abertas":  proc0.get("parcelas_abertas", ""),
        "parcelas_vencidas": proc0.get("parcelas_vencidas", ""),
        "existe_processo":   proc_existe(proc0),

        # Lista completa de processos (para Jinja {% for %})
        "processos": processos_ctx,

        # Cláusula 6
        "clausula_6_1": c61,
        "clausula_6_2": c62,

        # Quadro Resumo seção 3
        "custos_desc":   qr_custos["custos_desc"],
        "parcelamento":  qr_custos["parcelamento"],
        "n_parcelas_qr": qr_custos["n_parcelas_qr"],

        # Veículo (vazio se não for veículo)
        "veiculo_marca_modelo": vehicle.get("marca_modelo", ""),
        "veiculo_ano":          vehicle.get("ano", ""),
        "veiculo_placa":        vehicle.get("placa", ""),
        "veiculo_cor":          vehicle.get("cor", ""),
        "veiculo_renavam":      vehicle.get("renavam", ""),
        "veiculo_obs":          _build_veiculo_obs(vehicle),

        # Tipo do contrato
        "tipo_contrato": contract_type,

        # Cláusula 7.4 — taxa de rescisão antecipada
        "taxa_rescisao":         taxa_brl,
        "taxa_rescisao_extenso": taxa_extenso,
    }
    return ctx


# ── Mapeamento de template por tipo ──────────────────────────────────────────

TEMPLATE_KEYS = {
    "emprestimo":         "emprestimo",
    "veiculo":            "veiculo",
    "fiscal":             "fiscal",
    "condominio":         "condominio",
    "condominio_aluguel": "condominio_aluguel",
    "rural":              "rural",
}

# Tipos de contrato que recebem procuração extrajudicial automaticamente ao final
_EXTRAJUDICIAL_CONTRACT_TYPES = {"emprestimo", "veiculo", "fiscal", "rural"}


def _append_docx(base_path: Path, append_path: Path, output_path: Path) -> None:
    """Anexa o conteúdo de append_path ao final de base_path, salvando em output_path."""
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from copy import deepcopy

    base  = Document(str(base_path))
    extra = Document(str(append_path))

    # Quebra de página antes do conteúdo anexado
    p   = base.add_paragraph()
    run = p.add_run()
    br  = OxmlElement('w:br')
    br.set(qn('w:type'), 'page')
    run._r.append(br)

    # Copia todos os elementos do corpo, exceto sectPr final (evita conflito de formatação)
    for element in extra.element.body:
        if element.tag.endswith('}sectPr'):
            continue
        base.element.body.append(deepcopy(element))

    base.save(str(output_path))


def _docx_to_pdf(docx_path: Path) -> Path:
    """Converte DOCX para PDF e remove o DOCX."""
    pdf_path = docx_path.with_suffix(".pdf")
    try:
        from docx2pdf import convert
        convert(str(docx_path), str(pdf_path))
        docx_path.unlink()
        return pdf_path
    except Exception:
        return docx_path  # sem Word instalado: mantém o DOCX


def _pptx_to_pdf(pptx_path: Path) -> Path:
    """Converte PPTX → PDF via PowerPoint COM. Se não disponível, mantém .pptx."""
    pdf_path = pptx_path.with_suffix(".pdf")
    try:
        import win32com.client
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Visible = True
        prs = ppt.Presentations.Open(str(pptx_path.resolve()))
        prs.SaveAs(str(pdf_path.resolve()), 32)  # 32 = ppSaveAsPDF
        prs.Close()
        ppt.Quit()
        pptx_path.unlink()
        return pdf_path
    except Exception:
        return pptx_path  # sem PowerPoint: mantém o .pptx


# ── Proposta de Honorários ────────────────────────────────────────────────────

def _pptx_para_replace(para, old: str, new: str) -> bool:
    """Substitui old→new no parágrafo, consolidando todos os runs no primeiro."""
    full = "".join(r.text or "" for r in para.runs)
    if not full or old not in full:
        return False
    new_full = full.replace(old, new)
    if para.runs:
        para.runs[0].text = new_full
        for r in para.runs[1:]:
            r.text = ""
    return True


def _pptx_apply_shapes(shapes, replacements: list[tuple[str, str]]) -> None:
    """Aplica substituições em todas as shapes, recursivo para grupos."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _pptx_apply_shapes(shape.shapes, replacements)
        elif shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for old, new in replacements:
                    _pptx_para_replace(para, old, new)


def _brl2f(s: str) -> float:
    try:
        return float(s.strip().replace(".", "").replace(",", "."))
    except Exception:
        return 0.0


def _f2brl(v: float) -> str:
    return f"{v:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")


def build_proposta_replacements(payload: dict) -> list[tuple[str, str]]:
    """
    Retorna lista ordenada de (placeholder, valor_real) para preencher
    o template proposta.pptx.
    """
    nome    = payload.get("solicitante_nome",    "").strip()
    empresa = payload.get("solicitante_empresa", "").strip()
    assunto = payload.get("assunto",             "").strip()
    data    = payload.get("data_proposta",       "").strip()

    op1_vt  = payload.get("op1_valor_total",  "").strip()
    op1_ent = payload.get("op1_entrada",      "").strip()
    op1_np  = payload.get("op1_n_parcelas",   "").strip()
    op1_vp  = payload.get("op1_valor_parcela","").strip()

    op2_av  = payload.get("op2_avista",        "").strip()
    op2_np  = payload.get("op2_n_parcelas",    "").strip()
    op2_vp  = payload.get("op2_valor_parcela", "").strip()

    # Cálculos automáticos
    ec_val   = _brl2f(op1_vt) - _brl2f(op2_av)
    n2       = int(op2_np) if op2_np.isdigit() else 0
    tot_val  = n2 * _brl2f(op2_vp)

    ec_str   = _f2brl(ec_val)
    tot_str  = _f2brl(tot_val)

    # Partes envolvidas (slide 3): lista dinâmica separada por \n
    partes_raw: list = payload.get("partes_envolvidas", [])
    partes_str = "\n".join(p.strip() for p in partes_raw if p.strip()) or nome

    # Valor da causa e percentuais (slides 3, 5, 6, 7)
    valor_causa = payload.get("valor_causa", "").strip()
    causa_f = _brl2f(valor_causa) if valor_causa else 0.0
    if causa_f > 0:
        op1_f = _brl2f(op1_vt)
        op2_f = _brl2f(op2_av)
        perc_op1    = f"{op1_f / causa_f * 100:.1f}%"
        perc_op2_av = f"{op2_f / causa_f * 100:.1f}%"
    else:
        perc_op1    = "—"
        perc_op2_av = "—"

    return [
        # Strings mais longas/específicas primeiro
        ("{{OP2_ECONOMIA_LONGA}}",      f"Economia de R$ {ec_str} em relação ao parcelado por boleto. Aprovação rápida."),
        ("{{OP2_PARCELAS_SEM_JUROS}}", f"Até {op2_np} parcelas sem juros adicionais"),
        ("{{SOLICITANTE_EMPRESA}}",      empresa),
        ("{{ASSUNTO}}",                  assunto),
        ("{{SOLICITANTE_NOME}}",         nome),
        ("{{PARTES_LISTADAS}}",          partes_str),
        ("{{PARTES_S3_CLEAR}}",          ""),
        ("{{OP2_ECONOMIA_CURTA}}",      f"Economia de R$ {ec_str}"),
        ("{{OP2_TOTAL}}",               f"Total: R$ {tot_str}"),
        ("{{OP2_PARCELAS_MENSAIS}}",    f"em {op2_np} parcelas mensais"),
        ("{{OP1_PARCELAS}}",            f"{op1_np}x R$ {op1_vp}"),
        ("{{DATA_PROPOSTA}}",            data),
        ("{{OP2_ATE_PARCELAS}}",        f"Até {op2_np}x de"),
        ("{{PERC_OP2_AV}}",             perc_op2_av),
        ("{{PERC_OP1}}",                perc_op1),
        ("{{VALOR_CAUSA}}",             f"R$ {valor_causa}" if valor_causa else "—"),
        ("{{OP1_VALOR_TOTAL}}",         f"R$ {op1_vt}"),
        ("{{OP2_AVISTA}}",              f"R$ {op2_av}"),
        ("{{OP1_ENTRADA}}",             f"R$ {op1_ent}"),
        ("{{OP2_VALOR_PARCELA}}",       f"R$ {op2_vp}"),
    ]


def fill_proposta_template(template_path: Path, payload: dict, out_path: Path, to_pdf: bool = True) -> Path:
    """Preenche o template PPTX de proposta e opcionalmente converte para PDF."""
    from pptx import Presentation

    replacements = build_proposta_replacements(payload)
    prs = Presentation(str(template_path))

    for slide in prs.slides:
        _pptx_apply_shapes(slide.shapes, replacements)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return _pptx_to_pdf(out_path) if to_pdf else out_path


def generate_proposta(payload: dict, templates_dir: Path, output_dir: Path) -> tuple[list[Path], Path]:
    """Gera a proposta de honorários. Retorna (lista_arquivos, pasta_saída)."""
    tpl_path = templates_dir / "proposta.pptx"
    if not tpl_path.exists():
        raise FileNotFoundError("Template de Proposta não encontrado. Adicione-o na tela de Templates.")

    nome      = payload.get("solicitante_nome", "cliente").title()
    today_str = date.today().strftime("%Y%m%d")
    safe_name = "".join(c for c in nome if c.isalnum() or c in " _-").strip()[:40]
    to_pdf    = payload.get("formato", "pdf") == "pdf"
    ext       = "pdf" if to_pdf else "pptx"

    out_folder = output_dir / f"{safe_name} {today_str}"
    out_folder.mkdir(parents=True, exist_ok=True)

    out_file = fill_proposta_template(
        tpl_path, payload,
        out_folder / f"Proposta - {safe_name}.pptx",
        to_pdf,
    )
    # Renomear se converteu para PDF
    if to_pdf and out_file.suffix == ".pdf":
        final = out_folder / f"Proposta - {safe_name}.pdf"
        if out_file != final:
            out_file.rename(final)
            out_file = final

    return [out_file], out_folder


def fill_template(template_path: Path, context: dict, out_path: Path, to_pdf: bool = True) -> Path:
    """Preenche um template DOCX com docxtpl, salva e opcionalmente converte para PDF."""
    tpl = DocxTemplate(str(template_path))
    tpl.render(context)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    tpl.save(str(out_path))
    return _docx_to_pdf(out_path) if to_pdf else out_path


def generate_all(
    payload: dict,
    templates_dir: Path,
    output_dir: Path,
) -> list[Path]:
    """
    Gera todos os documentos selecionados.
    Retorna lista de caminhos gerados.
    """
    context      = build_context(payload)
    docs         = payload.get("docs", {})

    # Quando PJ, substitui bloco_cliente pela versão jurídica no contrato
    if docs.get("procuracao_pj", False):
        client  = payload.get("client", {})
        empresa = payload.get("empresa") or {}
        context["bloco_cliente"] = build_bloco_cliente_pj(client, empresa)

    contract_type = payload.get("contractType", "emprestimo")
    client_name  = payload.get("client", {}).get("nome", "cliente").title()
    to_pdf       = payload.get("formato", "pdf") == "pdf"

    # Pasta de saída: output_dir / nome_cliente_data
    today_str = date.today().strftime("%Y%m%d")
    safe_name = "".join(c for c in client_name if c.isalnum() or c in " _-").strip()[:40]
    out_folder = output_dir / f"{safe_name} {today_str}"
    out_folder.mkdir(parents=True, exist_ok=True)

    generated: list[Path] = []

    # Contrato
    if docs.get("contrato"):
        key = TEMPLATE_KEYS.get(contract_type, "emprestimo")
        tpl_path = templates_dir / f"{key}.docx"
        if not tpl_path.exists():
            raise FileNotFoundError(f"Template não encontrado: {tpl_path.name}\nAdicione-o na tela de Templates.")

        contract_docx = out_folder / f"Contrato - {safe_name}.docx"
        is_pj = docs.get("procuracao_pj", False)
        extrajud_tpl = templates_dir / (
            "procuracao_extrajudicial_pj.docx" if is_pj else "procuracao_extrajudicial.docx"
        )

        if contract_type in _EXTRAJUDICIAL_CONTRACT_TYPES and extrajud_tpl.exists():
            # Renderiza o contrato sem converter para PDF ainda
            fill_template(tpl_path, context, contract_docx, to_pdf=False)

            # Renderiza a procuração extrajudicial em arquivo temporário
            extrajud_tmp = out_folder / "_extrajudicial_tmp.docx"
            fill_template(extrajud_tpl, context, extrajud_tmp, to_pdf=False)

            # Anexa a extrajudicial ao final do contrato
            _append_docx(contract_docx, extrajud_tmp, contract_docx)
            extrajud_tmp.unlink(missing_ok=True)

            # Converte o documento unificado para PDF
            out = _docx_to_pdf(contract_docx) if to_pdf else contract_docx
        else:
            out = fill_template(tpl_path, context, contract_docx, to_pdf)

        generated.append(out)

    # Procuração
    if docs.get("procuracao"):
        is_pj    = docs.get("procuracao_pj", False)
        tpl_name = "procuracao_pj.docx" if is_pj else "procuracao.docx"
        tpl_path = templates_dir / tpl_name
        if not tpl_path.exists():
            raise FileNotFoundError(f"Template de Procuração {'PJ ' if is_pj else ''}não encontrado.")
        out = fill_template(tpl_path, context, out_folder / "Procuracao.docx", to_pdf)
        generated.append(out)

    # Hipossuficiência
    if docs.get("hipo"):
        tpl_path = templates_dir / "hipo.docx"
        if not tpl_path.exists():
            raise FileNotFoundError("Template de Hipossuficiência não encontrado.")
        out = fill_template(tpl_path, context, out_folder / "Hipossuficiencia.docx", to_pdf)
        generated.append(out)

    return generated, out_folder


def generate_contrato_veiculo(
    payload: dict,
    templates_dir: Path,
    output_dir: Path,
) -> tuple[Path, Path]:
    """Gera o Contrato de Compra de Veículo. Retorna (arquivo_gerado, pasta_saída)."""
    tpl_path = templates_dir / "contrato_veiculo.docx"
    if not tpl_path.exists():
        raise FileNotFoundError("Template contrato_veiculo.docx não encontrado.")

    nome    = payload.get("nome", "").upper()
    cidade  = payload.get("cidade", "Rio de Janeiro")
    uf      = payload.get("uf", "RJ")
    to_pdf  = payload.get("formato", "pdf") == "pdf"

    # Endereço completo em uma linha
    parts = [
        payload.get("rua", ""),
        payload.get("bairro", ""),
        payload.get("complemento", ""),
        f'{cidade}/{uf}',
        f'CEP: {payload.get("cep", "")}',
    ]
    endereco_completo = ", ".join(p for p in parts if p.strip())

    context = {
        "nome":               nome,
        "nacionalidade":      payload.get("nacionalidade", ""),
        "rg":                 payload.get("rg", ""),
        "cpf":                payload.get("cpf", ""),
        "endereco_completo":  endereco_completo,
        "email":              payload.get("email", ""),
        "veiculo_modelo":     payload.get("veiculo_modelo", "").upper(),
        "veiculo_ano":        payload.get("veiculo_ano", ""),
        "veiculo_cor":        payload.get("veiculo_cor", "").upper(),
        "veiculo_chassi":     payload.get("veiculo_chassi", "").upper(),
        "veiculo_placa":      payload.get("veiculo_placa", "").upper(),
        "veiculo_renavam":    payload.get("veiculo_renavam", ""),
        "local_data":         build_local_data(cidade, uf),
    }

    today_str = date.today().strftime("%Y%m%d")
    safe_name = "".join(c for c in nome.title() if c.isalnum() or c in " _-").strip()[:40]
    out_folder = output_dir / f"{safe_name} {today_str}"
    out_folder.mkdir(parents=True, exist_ok=True)

    out = fill_template(tpl_path, context, out_folder / "Contrato Veiculo.docx", to_pdf)
    return out, out_folder
